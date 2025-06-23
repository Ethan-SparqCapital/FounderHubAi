'''
Backend Code for Pitch Deck Generator App:

Next steps: update frontend to use new AI features (AI features in backend are omitted for now).
Test the app locally to ensure everything works. 
Deploy backend to a cloud platform such as AWS, Google Cloud, Azure, Heroku, etc. 
Deploy frontend on a static hosting service such as Vercel, Netlify, etc. 
'''

from fastapi import FastAPI, HTTPException, Depends, Body, Request # Importing the FastAPI class and the HTTPException class. 
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional, Dict, Literal, Any
from fastapi.responses import StreamingResponse
import io
import json
import os
import re

# Unomitted: 

import openai # Importing the openai library - to be implemented in the future. 
import os # Importing the os library - to be implemented in the future.
from dotenv import load_dotenv # Importing the dotenv library - to be implemented in the future. 
import time # Importing the time library - to be implemented in the future. 
from functools import lru_cache # Importing the lru_cache decorator - to be implemented in the future. - this is a decorator that caches the results of a function call based on the input arguments. 

# Load environment variables from .env file. 
load_dotenv()

# Set your OpenAI API key. 
openai.api_key = os.getenv("OPENAI_API_KEY") # This is the OpenAI API key in order to access the OpenAI API - this is a secret key that is stored in the .env file and is used to authenticate the user. 

# End of unomitted.

app = FastAPI()

# Add CORS middleware to allow requests from the frontend. 
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000"], # Need to eventually change this to production URL such as "https://pitchdeck.com" - your public domain.
    allow_credentials=True,
    allow_methods=["GET", "POST", "OPTIONS"],  # Explicitly allow OPTIONS (eventually change to "*" to allow all methods)
    allow_headers=["*"],
)

# ===== MODELS =====
class SlideRequest(BaseModel):  # creating a class called SlideRequest.
    problem: str # Create a new String for the object. 
    solution: str # Create a new String for the object. 

class SlideContentRequest(BaseModel):
    problem: str
    solution: str
    slide_title: str
    current_content: Optional[str] = None
    mode: Optional[str] = None  # New: mode can be 'optimize', 'improve', or None

class Slide(BaseModel):
    problem: str
    content: Optional[str] = None
    design: Optional[str] = None

# ===== Slide Analysis =====
class SlideAnalysisRequest(BaseModel):
    slides: List[Dict[str, str]]  # Each slide: {title, content}

class SlideAnalysisResponse(BaseModel):
    score: float
    narrative_flow: str
    visual_design: str
    data_credibility: str
    feedback: str

# ===== ENDPOINTS =====
# Define standard slide types
STANDARD_SLIDES: List[str] = [ # This is a list of standard slides which can be reused for both placeholder AND AI-generated implementation (adds "Literal" type for better safer type checking).
    "The Problem",
    "Our Solution",
    "Product Demo",
    "Market Opportunity",
    "Traction",
    "Customer Love",
    "Competitive Landscape",
    "Business Model",
    "Financial Projections",
    "Go-to-Market Strategy",
    "Team",
    "Funding Ask",
    "Thank You"
]

'''# Placeholder implementation (to be omitted) 
@app.post("/generate-slides")
async def generate_slides(request: SlideRequest):
    try:
        # Create slides with problem and solution in the first two slides
        slides = [
            f"The Problem: {request.problem}", # As a sample, we are adding the problem and solution inputs in the first two slides. 
            f"Our Solution: {request.solution}",
            "Product Demo",
            "Market Opportunity",
            "Traction",
            "Customer Love",
            "Competitive Landscape",
            "Business Model",
            "Financial Projections",
            "Go-to-Market Strategy",
            "Team",
            "Funding Ask",
            "Thank You"
        ]
        return {"slides": slides}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating slides: {str(e)}")
# End of block to omit. '''

# To be unomitted:

# This is a function that generates slides using OpenAI API. 
@app.post("/generate-slides")
async def generate_slides(request: SlideRequest):
    try:
        # Create a cache key from the request data - Cache Responses to save API costs
        cache_key = f"{request.problem}:{request.solution}"
        
        # Check if we have a cached response
        cached_response = get_cached_response(cache_key)
        if cached_response:
            return cached_response

        # Call OpenAI API to generate content for each slide. 
        client = openai.OpenAI()  # Initialize the client
        try:
            response = client.chat.completions.create(
                model="gpt-4", # Using model gpt-4 to generate the slides. 
            messages=[
                # This is the auto-generated system prompt sent to the OpenAI API behind the scenes, with the user's input for problem and solution being passed in along with the system prompt to generate the slides. 
                {"role": "system", "content": "You are a pitch deck expert. You are given a problem statement and a solution to the problem. You are to generate content for all slides in a startup pitch deck. For each slide, provide a compelling headline and 2-3 bullet points of key information. Make the content concise, impactful, and investor-ready, ensuring the slides are engaging and can instantly grab the attention of the audience and investors."},
                {"role": "user", "content": f"Generate content for all slides in a pitch deck about: Problem: '{request.problem}', Solution: '{request.solution}'. For each slide, provide a compelling headline and 2-3 bullet points of key information. The slides should be: The Problem, Our Solution, Product Demo, Market Opportunity, Traction, Customer Love, Competitive Landscape, Business Model, Financial Projections, Go-to-Market Strategy, Team, Funding Ask, and Thank You."}
            ],
            max_tokens=2000, # This is the maximum number of tokens that can be generated. 
        )
        except Exception as api_error:
            print(f"OpenAI API Error: {str(api_error)}")  # Print the error for debugging
            raise HTTPException(status_code=500, detail=f"OpenAI API error: {str(api_error)}")

        # Extract the generated content and split into slides
        if not response.choices or not response.choices[0].message:
            raise HTTPException(status_code=500, detail="No response content received from OpenAI")
            
        generated_text = response.choices[0].message.content
        if not generated_text:
            raise HTTPException(status_code=500, detail="Empty response content received from OpenAI")
            
        slides = [slide.strip() for slide in generated_text.split('\n') if slide.strip()]
        
        # Ensure we have all standard slides
        missing_slides = [slide for slide in STANDARD_SLIDES if not any(slide.lower() in s.lower() for s in slides)]
        if missing_slides:
            # Add any missing standard slides
            slides.extend(missing_slides)

        result = {"slides": slides}
        
        # Cache the response - Cache Responses to save API costs
        cache_response(cache_key, result)
        
        return result
    except HTTPException as he:
        # Re-raise HTTP exceptions
        raise he
    except Exception as e:
        # Log the full error for debugging
        print(f"Unexpected error in generate_slides: {str(e)}")
        print(f"Error type: {type(e)}")
        import traceback
        print(f"Traceback: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=f"An unexpected error occurred: {str(e)}")

# Simple in-memory cache implementation - Cache Responses to save API costs
_response_cache = {}
_cache_ttl = 3600  # Cache TTL in seconds (1 hour)

def get_cached_response(key: str):
    if key in _response_cache:
        timestamp, response = _response_cache[key]
        if time.time() - timestamp < _cache_ttl:
            return response
        else:
            del _response_cache[key]
    return None

def cache_response(key: str, response: dict):
    _response_cache[key] = (time.time(), response)

# ===== AI IMPLEMENTATION - SLIDE CONTENT ENDPOINT =====
@app.post("/generate-slide-content")
async def generate_slide_content(request: SlideContentRequest):
    try:
        # Validate slide title is in standard slides
        if request.slide_title not in STANDARD_SLIDES:
            raise HTTPException(status_code=400, detail=f"Invalid slide title. Must be one of: {', '.join(STANDARD_SLIDES)}")

        # Adjust prompt based on mode
        if request.mode == "optimize":
            prompt = f"""Given the pitch deck titled '{request.problem}' with description: '{request.solution}',\n"""
            if request.current_content:
                prompt += f"and the current content of the slide '{request.slide_title}': '{request.current_content}',\n"
            prompt += "Please optimize this slide to be more compelling and persuasive for investors. Focus on what investors care about most: market size, traction, defensibility, and growth potential. Provide a compelling headline and 2-3 bullet points of key information."
        elif request.mode == "improve":
            prompt = f"""Given the pitch deck titled '{request.problem}' with description: '{request.solution}',\n"""
            if request.current_content:
                prompt += f"and the current content of the slide '{request.slide_title}': '{request.current_content}',\n"
            prompt += "Please improve the messaging of this slide to be clearer, more persuasive, and more memorable. Provide a compelling headline and 2-3 bullet points of key information."
        else:
            # Default: generate or regenerate content
            if request.current_content:
                prompt = f"""Given the pitch deck titled '{request.problem}' with description: '{request.solution}',\n"""
                prompt += f"and the current content of the slide '{request.slide_title}': '{request.current_content}',\n"
                prompt += "please improve and enhance this slide's content while maintaining its core message. Make it more engaging and impactful for investors. Provide a compelling headline and 2-3 bullet points of key information."
            else:
                prompt = f"""Given the pitch deck titled '{request.problem}' with description: '{request.solution}',\n"""
                prompt += f"please generate detailed content for the slide '{request.slide_title}'. Make it engaging and impactful for investors. Provide a compelling headline and 2-3 bullet points of key information."

        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a pitch deck expert. Generate compelling and concise content for individual slides."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=2000,
        )

        content = response.choices[0].message.content.strip()
        return {"content": content}
    except openai.error.RateLimitError:
        raise HTTPException(status_code=429, detail="API rate limit exceeded. Please try again later.")
    except openai.error.APIError as e:
        raise HTTPException(status_code=500, detail=f"OpenAI API error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# ===== AI IMPLEMENTATION - DESIGN SUGGESTIONS ENDPOINT =====
@app.post("/generate-design-suggestions")
async def generate_design_suggestions(request: SlideContentRequest):
    try:
        # Validate slide title is in standard slides
        if request.slide_title not in STANDARD_SLIDES:
            raise HTTPException(status_code=400, detail=f"Invalid slide title. Must be one of: {', '.join(STANDARD_SLIDES)}")

        prompt = f"""Given the pitch deck titled '{request.problem}' with description: '{request.solution}',
        and the slide '{request.slide_title}' with content: '{request.current_content}',
        provide specific design suggestions to make this slide more visually appealing and effective.
        Include recommendations for:
        1. Layout and structure
        2. Visual elements (charts, images, icons)
        3. Color scheme
        4. Typography
        5. Data visualization (if applicable)"""

        response = openai.ChatCompletion.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are a presentation design expert. Provide specific and actionable design suggestions for pitch deck slides."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=500,
        )

        suggestions = response.choices[0].message.content.strip()
        return {"suggestions": suggestions}
    except openai.error.RateLimitError:
        raise HTTPException(status_code=429, detail="API rate limit exceeded. Please try again later.")
    except openai.error.APIError as e:
        raise HTTPException(status_code=500, detail=f"OpenAI API error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
# END OF AI IMPLEMENTATION 

#@app.get("/health")
#def health_check():
#    return {"status": "ok"}

# ===== SUGGESTION ENDPOINT =====
# Placeholder implementation for /generate-suggestion
@app.post("/generate-suggestion")
def generate_suggestion(
    data: Dict = Body(...)
):
    """
    Placeholder: Returns a static suggestion based on type.
    """
    suggestion_type = data.get("type", "Content")
    if suggestion_type == "Content":
        return {"suggestion": "Add a compelling statistic to the slide"}
    else:
        return {"suggestion": "Use a blue background with white text for emphasis"}
# End of block to omit.

# ===== AI IMPLEMENTATION - SUGGESTION ENDPOINT =====
# START OF AI IMPLEMENTATION - To be uncommented later
@app.post("/generate-suggestion")
async def generate_suggestion(
    data: Dict = Body(...)
):
    """
    AI implementation: Returns an AI-generated suggestion for the slide.
    """
    suggestion_type = data.get("type", "Content")
    slide_title = data.get("slide_title", "")
    content = data.get("content", "")
    design = data.get("design", "")
    if suggestion_type == "Content":
        prompt = f"""Given the slide titled '{slide_title}' with content: '{content}', suggest a single, actionable improvement to the slide's content for a startup pitch deck. Respond with only the suggestion."""
    else:
        prompt = f"""Given the slide titled '{slide_title}' with design notes: '{design}', suggest a single, actionable improvement to the slide's design (layout, visuals, colors, etc.) for a startup pitch deck. Respond with only the suggestion."""
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a pitch deck expert. Provide concise, actionable suggestions for improving slide content or design."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=100,
    )
    suggestion = response.choices[0].message.content.strip()
    return {"suggestion": suggestion}
#END OF AI IMPLEMENTATION 

# ===== VISUAL DATA ENDPOINT =====
# Placeholder implementation for /generate-visual-data
@app.post("/generate-visual-data")
def generate_visual_data(data: Dict = Body(...)):
    """
    Placeholder: Returns static chart/table data for the requested type.
    """
    visual_type = data.get("type", "pie")
    if visual_type == "pie":
        return {"data": [
            {"name": "A", "value": 40},
            {"name": "B", "value": 30},
            {"name": "C", "value": 30},
        ]}
    if visual_type == "bar":
        return {"data": [
            {"name": "Jan", "value": 30},
            {"name": "Feb", "value": 20},
            {"name": "Mar", "value": 50},
        ]}
    if visual_type == "line":
        return {"data": [
            {"name": "Q1", "value": 10},
            {"name": "Q2", "value": 40},
            {"name": "Q3", "value": 25},
        ]}
    if visual_type == "scatter":
        return {"data": [
            {"x": 10, "y": 20},
            {"x": 20, "y": 30},
            {"x": 30, "y": 10},
        ]}
    if visual_type == "table":
        return {"data": {
            "columns": ["Year", "Revenue", "Profit"],
            "rows": [
                ["2022", "$1M", "$200K"],
                ["2023", "$1.5M", "$350K"],
            ],
        }}
    return {"data": []}
# End of block to omit.

# ===== AI IMPLEMENTATION - VISUAL DATA ENDPOINT =====
# START OF AI IMPLEMENTATION - To be uncommented later
@app.post("/generate-visual-data")
async def generate_visual_data(data: Dict = Body(...)):
    """
    AI implementation: Returns chart/table data generated by OpenAI for the requested type and context.
    """
    visual_type = data.get("type", "pie")
    context = data.get("context", "")
    prompt = f"""Generate JSON data for a {visual_type} chart for a startup pitch deck. Context: {context}. 
    For pie/bar/line, use a list of objects with 'name' and 'value'. For scatter, use a list of objects with 'x' and 'y'. For table, use an object with 'columns' and 'rows'. Respond with only the JSON data."""
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are a data visualization expert. Generate chart/table data for pitch decks."},
            {"role": "user", "content": prompt}
        ],
        max_tokens=300,
    )
    import json
    try:
        data_json = json.loads(response.choices[0].message.content.strip())
    except Exception:
        data_json = response.choices[0].message.content.strip()
    return {"data": data_json}
# END OF AI IMPLEMENTATION 

# ===== PITCH DECK ANALYSIS ENDPOINT =====
@app.post("/analyze-pitch-deck", response_model=SlideAnalysisResponse)
async def analyze_pitch_deck(request: SlideAnalysisRequest):
    try:
        # Prepare the content for analysis
        slides_content = "\n\n".join([
            f"Slide: {slide.get('title', 'Untitled')}\nContent: {slide.get('content', '')}"
            for slide in request.slides
        ])

        # Call OpenAI API for analysis
        client = openai.OpenAI()
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": "You are an expert pitch deck reviewer. Analyze the pitch deck and provide detailed feedback on narrative flow, visual design, data credibility, and overall effectiveness. Provide specific, actionable suggestions for improvement."},
                {"role": "user", "content": f"Please analyze this pitch deck and provide feedback:\n\n{slides_content}\n\nProvide your analysis in the following distinct sections, using these exact titles, including the 'SECTION:' prefix:\n\nSECTION: Overall Score\n\nSECTION: Narrative Flow Analysis\n\nSECTION: Visual Design Analysis\n\nSECTION: Data Credibility Analysis\n\nSECTION: Specific Feedback and Suggestions"}
            ],
            max_tokens=2000,
        )

        # Parse the response
        analysis_text = response.choices[0].message.content.strip()
        
        parsed_data = {
            "score": 75.0,
            "narrative_flow": "No narrative flow analysis provided.",
            "visual_design": "No visual design analysis provided.",
            "data_credibility": "No data credibility analysis provided.",
            "feedback": "No specific feedback provided."
        }

        sections = analysis_text.split('SECTION:')[1:]

        for section in sections:
            section = section.strip()
            if 'Overall Score' in section:
                score_match = re.search(r'(\d+)', section)
                if score_match:
                    parsed_data["score"] = float(score_match.group(1))
            elif 'Narrative Flow Analysis' in section:
                parsed_data["narrative_flow"] = section.replace('Narrative Flow Analysis', '').strip()
            elif 'Visual Design Analysis' in section:
                parsed_data["visual_design"] = section.replace('Visual Design Analysis', '').strip()
            elif 'Data Credibility Analysis' in section:
                parsed_data["data_credibility"] = section.replace('Data Credibility Analysis', '').strip()
            elif 'Specific Feedback and Suggestions' in section:
                parsed_data["feedback"] = section.replace('Specific Feedback and Suggestions', '').strip()

        return SlideAnalysisResponse(
            score=parsed_data["score"],
            narrative_flow=parsed_data["narrative_flow"],
            visual_design=parsed_data["visual_design"],
            data_credibility=parsed_data["data_credibility"],
            feedback=parsed_data["feedback"]
        )
    except Exception as e:
        print(f"Error in analyze_pitch_deck: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Error analyzing pitch deck: {str(e)}")

# PDF export (requires reportlab)
@app.post("/export-pdf")
def export_pdf(request: SlideAnalysisRequest):
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        buffer = io.BytesIO()
        p = canvas.Canvas(buffer, pagesize=letter)
        width, height = letter
        y = height - 50
        for slide in request.slides:
            p.setFont("Helvetica-Bold", 16)
            p.drawString(50, y, slide['title'])
            y -= 30
            p.setFont("Helvetica", 12)
            for line in (slide.get('content', '') or '').split('\n'):
                p.drawString(60, y, line)
                y -= 20
                if y < 80:
                    p.showPage()
                    y = height - 50
            y -= 20
            if y < 80:
                p.showPage()
                y = height - 50
        p.save()
        buffer.seek(0)
        return StreamingResponse(buffer, media_type="application/pdf", headers={"Content-Disposition": "attachment; filename=pitch_deck.pdf"})
    except ImportError:
        raise HTTPException(status_code=500, detail="reportlab is required for PDF export. Please install it.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# PPTX export (requires python-pptx)
@app.post("/export-ppt")
def export_ppt(request: SlideAnalysisRequest):
    try:
        from pptx import Presentation
        prs = Presentation()
        for slide_data in request.slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = slide_data['title']
            content = slide.placeholders[1]
            content.text = slide_data.get('content', '') or ''
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return StreamingResponse(buffer, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", headers={"Content-Disposition": "attachment; filename=pitch_deck.pptx"})
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx is required for PPTX export. Please install it.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

# Requirements:
# pip install reportlab python-pptx

DATA_FILE = "user_data.json"

def load_user_data(user_id):
    if not os.path.exists(DATA_FILE):
        return {}
    with open(DATA_FILE, "r") as f:
        data = json.load(f)
    return data.get(user_id, {})

def save_user_data(user_id, user_data):
    if os.path.exists(DATA_FILE):
        with open(DATA_FILE, "r") as f:
            data = json.load(f)
    else:
        data = {}
    data[user_id] = user_data
    with open(DATA_FILE, "w") as f:
        json.dump(data, f)

@app.post("/save-slides")
async def save_slides(request: Request):
    body = await request.json()
    user_id = body.get("userId", "demo")  # Replace with real user/session ID
    slides = body.get("slides", [])
    user_data = load_user_data(user_id)
    user_data["slides"] = slides
    save_user_data(user_id, user_data)
    return {"status": "ok"}

@app.get("/get-slides")
async def get_slides(userId: str = "demo"):
    user_data = load_user_data(userId)
    return {"slides": user_data.get("slides", [])}

@app.get("/dashboard-stats")
async def dashboard_stats(userId: str = "demo"):
    user_data = load_user_data(userId)
    slides = user_data.get("slides", [])
    # You can add more stats as you add more features
    return {
        "decksCreated": 1 if slides else 0,
        "legalDocsGenerated": user_data.get("legalDocsGenerated", 0),
        "researchReports": user_data.get("researchReports", 0),
    }